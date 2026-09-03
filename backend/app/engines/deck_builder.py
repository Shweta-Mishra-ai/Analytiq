"""
engines/deck_builder.py — the report as a deck.

In the consulting world Analytiq is benchmarked against, the deck is the
deliverable. A PDF report gets read once by whoever commissioned it; a
slide gets presented to the people who decide. Everything here is already
computed for the PDF — the KPIs, the findings, the decision bands, the
charts as PNG bytes — so this is a second rendering of the same analysis,
not a second analysis.

Typography deliberately mirrors the PDF: that report is set in Carlito
and Caladea, which are metric-compatible with Calibri and Cambria. Using
those here means a deck and a report from the same run look like two
parts of one document rather than two tools' output.
"""
from __future__ import annotations

import io
import logging
from datetime import datetime
from typing import Dict, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# 16:9 at the size PowerPoint opens by default.
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)
CONTENT_W = SLIDE_W - 2 * MARGIN

FONT_HEAD = "Cambria"     # metric-compatible with the report's Caladea
FONT_BODY = "Calibri"     # metric-compatible with the report's Carlito


def _rgb(hex_str: str) -> RGBColor:
    h = str(hex_str or "#000000").lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    try:
        return RGBColor.from_string(h.upper()[:6])
    except Exception:
        logger.debug("bad colour %r — using black", hex_str)
        return RGBColor(0, 0, 0)


class DeckTheme:
    """Colours pulled from the same theme the PDF used, so the two match."""

    def __init__(self, t: Dict):
        self.ink = _rgb(t.get("text", "#1A1F28"))
        self.muted = _rgb(t.get("text_muted", "#6B7585"))
        self.accent = _rgb(t.get("accent", "#14555F"))
        self.dark = _rgb(t.get("cover_bg", "#0E1B2A"))
        self.light = _rgb("#FFFFFF")
        self.panel = _rgb(t.get("bg_light", "#F1F4F8"))
        self.positive = _rgb(t.get("positive", "#256B47"))
        self.warning = _rgb(t.get("warning", "#8A5B0E"))
        self.negative = _rgb(t.get("negative", "#A32C1C"))
        self.border = _rgb(t.get("border", "#D9DEE6"))

    def severity(self, sev: str) -> RGBColor:
        return {
            "critical": self.negative, "high": self.warning,
            "warning": self.warning, "medium": self.warning,
            "positive": self.positive,
        }.get(str(sev).lower(), self.accent)


# ══════════════════════════════════════════════════════════
#  PRIMITIVES
# ══════════════════════════════════════════════════════════

def _blank(prs: Presentation):
    """A slide with no placeholders — every element is positioned here, so
    an inherited layout cannot move something underneath us."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill(slide, colour: RGBColor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colour


def _text(slide, left, top, width, height, runs, *, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, spacing=None):
    """A text box built from (text, size, bold, colour, font) runs.

    Built run by run rather than through `text_frame.text`, which collapses
    a paragraph to one unstyled run and loses every distinction inside it.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    first = True
    for text, size, bold, colour, font in runs:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = align
        if spacing is not None:
            para.space_after = spacing
        run = para.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour
        run.font.name = font
    return box


def _card(slide, left, top, width, height, theme: DeckTheme,
          fill: Optional[RGBColor] = None):
    """A tinted block. Deliberately no edge stripe — an accent bar down one
    side of a card is the most recognisable tell of generated slides."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or theme.panel
    shape.line.fill.background()
    shape.shadow.inherit = False
    try:
        shape.adjustments[0] = 0.04
    except Exception:
        logger.debug("could not set corner radius", exc_info=True)
    return shape


def _slide_title(slide, theme: DeckTheme, title: str, kicker: str = "",
                 on_dark: bool = False):
    ink = theme.light if on_dark else theme.ink
    muted = theme.light if on_dark else theme.muted
    top = MARGIN
    if kicker:
        _text(slide, MARGIN, top, CONTENT_W, Inches(0.3),
              [(kicker.upper(), 11, True, theme.accent if not on_dark
                else theme.light, FONT_BODY)])
        top = top + Inches(0.34)
    _text(slide, MARGIN, top, CONTENT_W, Inches(0.8),
          [(title, 34, True, ink, FONT_HEAD)])
    return top + Inches(0.95)


# ══════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════

def _title_slide(prs, theme, config, domain_label, n_rows, n_cols):
    slide = _blank(prs)
    _fill(slide, theme.dark)
    _text(slide, MARGIN, Inches(2.0), CONTENT_W, Inches(0.4),
          [(domain_label.upper(), 12, True, theme.light, FONT_BODY)])
    _text(slide, MARGIN, Inches(2.5), CONTENT_W, Inches(1.6),
          [(config.get("title", "Data Analysis Report"), 48, True,
            theme.light, FONT_HEAD)])
    subtitle = config.get("subtitle") or ""
    if subtitle:
        _text(slide, MARGIN, Inches(4.1), CONTENT_W, Inches(0.5),
              [(subtitle, 16, False, theme.light, FONT_BODY)])
    _text(slide, MARGIN, Inches(5.9), CONTENT_W, Inches(1.0),
          [("Prepared for {}".format(config.get("client_name", "Client")),
            14, True, theme.light, FONT_BODY),
           ("{}  ·  {:,} records × {} fields".format(
               datetime.now().strftime("%d %B %Y"), n_rows, n_cols),
            11, False, theme.light, FONT_BODY)],
          spacing=Pt(4))
    slide.notes_slide.notes_text_frame.text = (
        "Generated by Analytiq from the submitted dataset. Every figure in "
        "this deck is computed from that data.")
    return slide


def _kpi_slide(prs, theme, kpis: Sequence, domain_label: str):
    """Large stat callouts. The numbers a reader of this kind of data looks
    for first, sized so they read from the back of a room."""
    slide = _blank(prs)
    _fill(slide, theme.light)
    top = _slide_title(slide, theme, "At a glance", domain_label)
    cards = list(kpis)[:4]
    if not cards:
        return slide
    gap = Inches(0.28)
    width = int((CONTENT_W - gap * (len(cards) - 1)) / len(cards))
    for i, k in enumerate(cards):
        left = MARGIN + i * (width + gap)
        _card(slide, left, top, width, Inches(2.1), theme)
        value = k.get("value")
        unit = k.get("unit") or ""
        shown = ("—" if value is None
                 else "{:,.0f}{}".format(value, unit)
                 if abs(value) >= 100 else "{:,.2f}{}".format(value, unit))
        pad = Inches(0.28)
        _text(slide, left + pad, top + Inches(0.3),
              width - 2 * pad, Inches(0.35),
              [(str(k.get("label", "")).upper(), 10, True, theme.muted,
                FONT_BODY)])
        _text(slide, left + pad, top + Inches(0.72),
              width - 2 * pad, Inches(0.8),
              [(shown, 38, True, theme.accent, FONT_HEAD)])
        note = k.get("benchmark") or k.get("note") or ""
        if note:
            _text(slide, left + pad, top + Inches(1.5),
                  width - 2 * pad, Inches(0.5),
                  [(note[:90], 8, False, theme.muted, FONT_BODY)])
    return slide


def _summary_slide(prs, theme, summary: str, findings: Sequence[str]):
    slide = _blank(prs)
    _fill(slide, theme.light)
    top = _slide_title(slide, theme, "Executive summary")
    col = int(CONTENT_W * 0.56)
    _text(slide, MARGIN, top, col, Inches(3.6),
          [(summary or "Analysis completed.", 15, False, theme.ink,
            FONT_BODY)], spacing=Pt(8))
    right_left = MARGIN + col + Inches(0.5)
    right_w = CONTENT_W - col - Inches(0.5)
    if findings:
        _card(slide, right_left, top - Inches(0.15), right_w,
              Inches(3.9), theme)
        runs = [("KEY FINDINGS", 10, True, theme.muted, FONT_BODY)]
        for f in list(findings)[:5]:
            runs.append((str(f)[:150], 11, False, theme.ink, FONT_BODY))
        _text(slide, right_left + Inches(0.3), top + Inches(0.15),
              right_w - Inches(0.6), Inches(3.5), runs, spacing=Pt(7))
    return slide


def _findings_slides(prs, theme, insights: Sequence):
    """Two findings per slide, each with its severity, evidence and first
    action. More than two and nothing on the slide reads from a distance."""
    made = []
    items = list(insights)[:6]
    for chunk_start in range(0, len(items), 2):
        chunk = items[chunk_start:chunk_start + 2]
        slide = _blank(prs)
        _fill(slide, theme.light)
        top = _slide_title(
            slide, theme,
            "What the data shows" if chunk_start == 0 else "What the data shows (continued)")
        for j, ins in enumerate(chunk):
            block_top = top + j * Inches(2.35)
            sev = str(getattr(ins, "severity", "info"))
            _card(slide, MARGIN, block_top, CONTENT_W, Inches(2.05), theme)
            pad = Inches(0.35)
            _text(slide, MARGIN + pad, block_top + Inches(0.22),
                  Inches(1.5), Inches(0.3),
                  [(sev.upper(), 9, True, theme.severity(sev), FONT_BODY)])
            _text(slide, MARGIN + pad, block_top + Inches(0.55),
                  CONTENT_W - 2 * pad, Inches(0.5),
                  [(str(getattr(ins, "title", ""))[:110], 18, True,
                    theme.ink, FONT_HEAD)])
            action = str(getattr(ins, "action", "")).split("2.")[0]
            action = action.lstrip("1.").strip()
            _text(slide, MARGIN + pad, block_top + Inches(1.12),
                  CONTENT_W - 2 * pad, Inches(0.8),
                  [(str(getattr(ins, "evidence", ""))[:150], 11, False,
                    theme.muted, FONT_BODY),
                   ("Do next: " + action[:120], 11, True, theme.ink,
                    FONT_BODY)], spacing=Pt(4))
        made.append(slide)
    return made


def _decision_slide(prs, theme, bands: Sequence, target: str):
    """The slide that turns a model into a decision."""
    slide = _blank(prs)
    _fill(slide, theme.light)
    top = _slide_title(slide, theme, "Where to act",
                       "predicting {}".format(str(target).replace("_", " ")))
    rows = list(bands)[:4]
    if not rows:
        return slide
    headers = ["If you act on", "Records", "Events reached", "Hit rate",
               "Share of events", "vs random"]
    table_h = Inches(0.45) * (len(rows) + 1)
    shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), MARGIN, top, CONTENT_W, table_h)
    table = shape.table
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.text_frame.paragraphs[0].add_run().text = head
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size, run.font.bold = Pt(11), True
        run.font.color.rgb = theme.light
        run.font.name = FONT_BODY
        cell.fill.solid()
        cell.fill.fore_color.rgb = theme.accent
    for r, b in enumerate(rows, 1):
        values = [
            "top {}%".format(b.budget_pct),
            "{:,}".format(b.n_targeted),
            "{:,} of {:,}".format(b.n_events_caught, b.total_events),
            "{:.0f}%".format(b.precision),
            "{:.0f}%".format(b.recall),
            "{:.1f}x".format(b.lift),
        ]
        for c, value in enumerate(values):
            cell = table.cell(r, c)
            cell.text_frame.paragraphs[0].add_run().text = value
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(12)
            run.font.bold = (c == 5)
            run.font.color.rgb = theme.ink
            run.font.name = FONT_BODY
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                theme.panel if r % 2 else theme.light)
    best = max(rows, key=lambda b: b.lift)
    _text(slide, MARGIN, top + table_h + Inches(0.4), CONTENT_W, Inches(1.0),
          [("Targeting the top {}% reaches {:,} of {:,} events, and {:.0f}% "
            "of those contacted record it — {:.1f} times better than "
            "choosing at random.".format(
                best.budget_pct, best.n_events_caught, best.total_events,
                best.precision, best.lift), 14, False, theme.ink, FONT_BODY),
           ("Which row to choose is a budget decision, not a modelling one.",
            12, True, theme.accent, FONT_BODY)], spacing=Pt(6))
    return slide


def _chart_slides(prs, theme, chart_data: Sequence[Tuple]):
    """One chart per slide, with its takeaway beside it rather than under
    it — a caption below a full-width image gets read last or not at all."""
    made = []
    for entry in list(chart_data)[:6]:
        # Entries may carry a ChartSpec after the narrative;
        # the deck does not need it, but must not choke on it.
        title, img_bytes, narrative = entry[0], entry[1], entry[2]
        if not img_bytes:
            continue
        slide = _blank(prs)
        _fill(slide, theme.light)
        top = _slide_title(slide, theme, str(title)[:70], "exhibit")
        # Fit inside the available box on BOTH axes. Scaling by width
        # alone pushed a square correlation matrix off the bottom of the
        # slide — charts here vary from wide bars to square heatmaps.
        box_w = int(CONTENT_W * 0.62)
        box_h = int(SLIDE_H - top - MARGIN)
        try:
            with io.BytesIO(img_bytes) as probe:
                from PIL import Image as _PILImage
                with _PILImage.open(probe) as im:
                    px_w, px_h = im.size
            ratio = px_h / px_w if px_w else 0.6
            draw_w = box_w
            draw_h = int(draw_w * ratio)
            if draw_h > box_h:
                draw_h = box_h
                draw_w = int(draw_h / ratio) if ratio else box_w
            slide.shapes.add_picture(io.BytesIO(img_bytes), MARGIN, top,
                                     width=draw_w, height=draw_h)
        except Exception:
            logger.warning("could not place chart %r", title, exc_info=True)
            continue
        if narrative:
            left = MARGIN + box_w + Inches(0.45)
            _text(slide, left, top, CONTENT_W - box_w - Inches(0.45),
                  Inches(3.8),
                  [("WHAT THIS SHOWS", 10, True, theme.muted, FONT_BODY),
                   (str(narrative)[:520], 12, False, theme.ink, FONT_BODY)],
                  spacing=Pt(8))
        made.append(slide)
    return made


def _outlook_slide(prs, theme, fc):
    """Where the measure is heading, or why no projection is shown.

    Both are worth a slide. "This is a level with noise around it" stops a
    room arguing about a trend that is not there.
    """
    slide = _blank(prs)
    _fill(slide, theme.light)
    measure = str(fc.column).replace("_", " ")
    top = _slide_title(slide, theme, "Outlook", measure)

    if not fc.usable:
        _card(slide, MARGIN, top, CONTENT_W, Inches(2.2), theme)
        _text(slide, MARGIN + Inches(0.35), top + Inches(0.3),
              CONTENT_W - Inches(0.7), Inches(1.6),
              [("No forecast shown", 20, True, theme.ink, FONT_HEAD),
               (fc.verdict[:340], 13, False, theme.ink, FONT_BODY)],
              spacing=Pt(8))
        return slide

    points = list(fc.points)[:4]
    gap = Inches(0.28)
    width = int((CONTENT_W - gap * (len(points) - 1)) / max(len(points), 1))
    for i, p in enumerate(points):
        left = MARGIN + i * (width + gap)
        _card(slide, left, top, width, Inches(1.9), theme)
        pad = Inches(0.25)
        _text(slide, left + pad, top + Inches(0.25), width - 2 * pad,
              Inches(0.3), [(p.period, 10, True, theme.muted, FONT_BODY)])
        _text(slide, left + pad, top + Inches(0.62), width - 2 * pad,
              Inches(0.6),
              [("{:,.0f}".format(p.value), 30, True, theme.accent,
                FONT_HEAD)])
        _text(slide, left + pad, top + Inches(1.3), width - 2 * pad,
              Inches(0.4),
              [("{:,.0f} – {:,.0f}".format(p.lower, p.upper), 10, False,
                theme.muted, FONT_BODY)])

    _text(slide, MARGIN, top + Inches(2.3), CONTENT_W, Inches(1.6),
          [(fc.verdict[:320], 14, False, theme.ink, FONT_BODY),
           ("The interval widens with distance because uncertainty "
            "compounds. Plan against the range, not the line.",
            12, True, theme.accent, FONT_BODY)], spacing=Pt(8))
    return slide


def _actions_slide(prs, theme, actions: Sequence[str]):
    slide = _blank(prs)
    _fill(slide, theme.light)
    top = _slide_title(slide, theme, "Recommended actions")
    items = list(actions)[:5]
    if not items:
        _text(slide, MARGIN, top, CONTENT_W, Inches(1.0),
              [("No action met the evidence threshold for inclusion.", 14,
                False, theme.muted, FONT_BODY)])
        return slide
    for i, action in enumerate(items):
        row_top = top + i * Inches(0.95)
        _card(slide, MARGIN, row_top, CONTENT_W, Inches(0.8), theme)
        _text(slide, MARGIN + Inches(0.3), row_top + Inches(0.2),
              Inches(0.5), Inches(0.4),
              [(str(i + 1), 20, True, theme.accent, FONT_HEAD)])
        _text(slide, MARGIN + Inches(0.95), row_top + Inches(0.22),
              CONTENT_W - Inches(1.4), Inches(0.5),
              [(str(action)[:170], 13, False, theme.ink, FONT_BODY)])
    return slide


def _method_slide(prs, theme, limitations: Sequence[str]):
    slide = _blank(prs)
    _fill(slide, theme.dark)
    _text(slide, MARGIN, Inches(1.4), CONTENT_W, Inches(0.9),
          [("How to read this", 34, True, theme.light, FONT_HEAD)])
    runs = []
    for line in limitations:
        runs.append((str(line), 13, False, theme.light, FONT_BODY))
    _text(slide, MARGIN, Inches(2.6), int(CONTENT_W * 0.8), Inches(3.6),
          runs, spacing=Pt(10))
    return slide


# ══════════════════════════════════════════════════════════
#  ENTRY POINT
# ══════════════════════════════════════════════════════════

def build_deck(
    df,
    config: dict,
    *,
    domain: str = "general",
    kpis: Sequence = (),
    executive_summary: str = "",
    findings: Sequence[str] = (),
    top_insights: Sequence = (),
    recommendations: Sequence[str] = (),
    chart_data: Sequence[Tuple] = (),
    predictive=None,
    forecast=None,
) -> bytes:
    """The same analysis as the PDF, as a deck.

    Takes what the report already computed rather than recomputing it, so
    a deck and a report generated from one dataset cannot disagree.
    """
    from app.engines.pdf.theme import THEMES
    from app.engines.pdf_builder import _domain_label, _domain_theme

    theme_name = config.get("theme_name")
    if theme_name not in THEMES:
        theme_name = _domain_theme(domain)
    theme = DeckTheme(THEMES[theme_name])
    domain_label = _domain_label(domain)

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H

    n_rows = len(df) if df is not None else 0
    n_cols = len(df.columns) if df is not None else 0

    _title_slide(prs, theme, config, domain_label, n_rows, n_cols)
    if kpis:
        _kpi_slide(prs, theme, kpis, domain_label)
    _summary_slide(prs, theme, executive_summary, findings)
    if top_insights:
        _findings_slides(prs, theme, top_insights)

    bands = list(getattr(predictive, "decision_bands", None) or [])
    if bands:
        _decision_slide(prs, theme, bands, getattr(predictive, "target", ""))

    if forecast is not None:
        _outlook_slide(prs, theme, forecast)

    if chart_data:
        _chart_slides(prs, theme, chart_data)

    _actions_slide(prs, theme, recommendations)

    limits = [
        "Every figure is computed from the submitted dataset. Nothing is "
        "estimated from outside it.",
        "Where a published range is cited it is general industry guidance, "
        "not a licensed benchmark, and it moves with sector and region.",
        "Findings shown here cleared both statistical significance and a "
        "minimum effect size — a difference too small to act on is not "
        "reported however certain it is.",
    ]
    if predictive is not None and getattr(predictive, "model_choice", None):
        choice = predictive.model_choice
        limits.append(
            "The model was selected from {} candidates on cross-validated "
            "ranking quality; its operating threshold is {:.2f}, chosen "
            "rather than assumed.".format(
                len(choice.candidates or []), choice.threshold))
    _method_slide(prs, theme, limits)

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out.read()
