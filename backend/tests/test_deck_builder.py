"""
The report as a deck.

In consulting the deck is the deliverable — a PDF gets read once by
whoever commissioned it, a slide gets presented to the people who decide.
This renders the analysis the PDF already computed, so the two cannot
disagree about the same dataset.

Rendering is unavailable in CI, so the checks visual inspection would do
— overflow, out-of-bounds shapes, overlaps — are made from shape geometry.
"""
import io

import numpy as np
import pandas as pd
import pytest
from pptx import Presentation
from pptx.util import Inches

from app.engines.deck_builder import SLIDE_H, SLIDE_W, build_deck

CONFIG = {"title": "Workforce Review", "client_name": "Acme Ltd",
          "subtitle": "Quarterly", "confidential": True, "theme_name": "",
          "logo_path": None, "prepared_by": "", "source_table": "src"}


@pytest.fixture(scope="module")
def inputs():
    from app.engines.chart_exporter import generate_all_charts
    from app.engines.kpi_engine import compute_kpis
    from app.engines.predictive import compute_drivers, find_binary_target
    from app.engines.story_engine import generate_story

    rng = np.random.default_rng(11)
    n = 900
    ot = rng.choice(["Yes", "No"], n, p=[.35, .65])
    ten = rng.integers(1, 20, n)
    p = (.05 + .32 * (ot == "Yes") + .24 * (ten <= 2)).clip(0, .92)
    df = pd.DataFrame({
        "employee_id": np.arange(n),
        "Attrition": np.where(rng.random(n) < p, "Yes", "No"),
        "OverTime": ot, "YearsAtCompany": ten,
        "JobRole": rng.choice(["Sales Rep", "Engineer", "Manager"], n),
        "MonthlyIncome": rng.uniform(2500, 19000, n).round(2),
    })
    story = generate_story(df)
    return {
        "df": df,
        "kpis": [c.as_dict() for c in compute_kpis(df, "hr")],
        "executive_summary": story.executive_summary,
        "findings": story.key_findings,
        "top_insights": story.top_insights,
        "recommendations": story.recommended_actions,
        "chart_data": [(t, b, "What this chart shows, in one paragraph.")
                       for t, b in generate_all_charts(df, "HR Blue", 3) if b],
        "predictive": compute_drivers(df, find_binary_target(df)),
    }


@pytest.fixture(scope="module")
def deck(inputs):
    data = dict(inputs)
    df = data.pop("df")
    return Presentation(io.BytesIO(
        build_deck(df=df, config=dict(CONFIG), domain="hr", **data)))


# ── it is a real deck ─────────────────────────────────────

def test_it_produces_a_valid_presentation(deck):
    assert len(deck.slides) >= 6


def test_it_opens_at_sixteen_by_nine(deck):
    """PowerPoint's default canvas is 10 inches wide. Coordinates past the
    edge are written rather than clamped, so a wrong canvas silently drops
    content off the slide."""
    assert deck.slide_width == SLIDE_W
    assert deck.slide_height == SLIDE_H


# ── it says the right things ──────────────────────────────

def _all_text(deck):
    """Every string in the deck, including table cells.

    Table text lives inside a graphic frame rather than on the shape's own
    text frame, so walking `has_text_frame` alone misses every header and
    cell — which is most of the decision table.
    """
    parts = []
    for slide in deck.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                parts.append(sh.text_frame.text)
            if getattr(sh, "has_table", False):
                for row in sh.table.rows:
                    for cell in row.cells:
                        parts.append(cell.text_frame.text)
    return "\n".join(parts)


def test_the_deck_carries_the_client_and_title(deck):
    text = _all_text(deck)
    assert "Workforce Review" in text
    assert "Acme Ltd" in text


def test_it_leads_with_the_numbers_then_the_argument(deck):
    """A deck that opens with methodology loses the room."""
    titles = [sh.text_frame.text for slide in deck.slides
              for sh in slide.shapes if sh.has_text_frame]
    joined = " | ".join(titles)
    assert joined.index("At a glance") < joined.index("Executive summary")


def test_the_decision_table_reaches_the_deck(deck, inputs):
    if not getattr(inputs["predictive"], "decision_bands", None):
        pytest.skip("no decision bands for this fixture")
    text = _all_text(deck)
    assert "Where to act" in text
    assert "vs random" in text
    assert "budget decision" in text


def test_charts_are_placed_with_their_takeaway(deck):
    pictures = [sh for slide in deck.slides for sh in slide.shapes
                if sh.shape_type is not None and sh.has_text_frame is False]
    assert pictures, "no chart images placed"
    assert "WHAT THIS SHOWS" in _all_text(deck)


def test_it_states_its_own_limits(deck):
    """The same basis-of-preparation the report carries."""
    text = _all_text(deck)
    assert "computed from the submitted dataset" in text


def test_speaker_notes_are_notes_not_a_text_box(deck):
    first = deck.slides[0]
    assert first.has_notes_slide
    assert first.notes_slide.notes_text_frame.text.strip()


# ── geometry, in place of a render ────────────────────────

def test_nothing_falls_off_the_slide(deck):
    for n, slide in enumerate(deck.slides, 1):
        for sh in slide.shapes:
            if sh.left is None or sh.top is None:
                continue
            assert sh.left >= 0 and sh.top >= 0, f"slide {n}: negative origin"
            assert sh.left + (sh.width or 0) <= SLIDE_W, \
                f"slide {n}: '{sh.shape_id}' runs past the right edge"
            assert sh.top + (sh.height or 0) <= SLIDE_H, \
                f"slide {n}: '{sh.shape_id}' runs past the bottom edge"


def test_a_square_chart_is_scaled_to_fit():
    """Scaling by width alone pushed a square correlation matrix off the
    bottom of the slide; charts range from wide bars to square heatmaps."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig, ax = plt.subplots(figsize=(6, 6))       # deliberately square
    ax.imshow(np.random.default_rng(1).random((8, 8)))
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=100)
    plt.close(fig)

    df = pd.DataFrame({"a": range(50), "b": range(50)})
    pptx = build_deck(df=df, config=dict(CONFIG), domain="general",
                      chart_data=[("Square chart", buf.getvalue(), "Note.")])
    prs = Presentation(io.BytesIO(pptx))
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.left is None:
                continue
            assert sh.top + (sh.height or 0) <= SLIDE_H


def test_text_boxes_do_not_overlap_each_other(deck):
    for n, slide in enumerate(deck.slides, 1):
        boxes = [(sh.left, sh.top, sh.width, sh.height)
                 for sh in slide.shapes
                 if sh.has_text_frame and sh.text_frame.text.strip()
                 and sh.left is not None]
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                a, b = boxes[i], boxes[j]
                ox = min(a[0] + a[2], b[0] + b[2]) - max(a[0], b[0])
                oy = min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1])
                assert not (ox > Inches(0.1) and oy > Inches(0.1)), \
                    f"slide {n}: two text boxes overlap"


# ── degenerate input ──────────────────────────────────────

def test_a_deck_builds_from_almost_nothing():
    """A dataset with no findings, no charts and no model must still
    produce something openable rather than raising."""
    df = pd.DataFrame({"a": [1, 2, 3]})
    data = build_deck(df=df, config=dict(CONFIG), domain="general")
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) >= 3


def test_the_endpoint_returns_a_presentation(client, uploaded_dataset_id):
    r = client.post(f"/api/reports/{uploaded_dataset_id}/pdf",
                    json={"title": "Deck", "format": "pptx", "max_charts": 2})
    assert r.status_code == 200, r.text
    assert "presentationml" in r.headers["content-type"]
    assert r.content[:2] == b"PK"          # it is a real zip container


def test_pdf_is_still_the_default(client, uploaded_dataset_id):
    r = client.post(f"/api/reports/{uploaded_dataset_id}/pdf",
                    json={"title": "Doc", "max_charts": 1})
    assert r.status_code == 200
    assert r.headers["content-type"] == "application/pdf"
